"""Extract plain text from common file types (PDF, Office, etc.)."""

from __future__ import annotations

import base64
import csv
import io
from pathlib import Path

_IMAGE_EXT = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}


def pdf_to_images_b64(data: bytes, max_pages: int = 2) -> list[str]:
    """
    Render PDF pages to PNG base64 for vision models (OCR via vision).

    Uses PyMuPDF (fitz). If PyMuPDF isn't installed or rendering fails, returns [].
    """
    if not data:
        return []
    try:
        import fitz  # type: ignore

        doc = fitz.open(stream=data, filetype="pdf")
        images: list[str] = []
        for i in range(min(max_pages, len(doc))):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=200)
            img_bytes = pix.tobytes("png")
            images.append(base64.b64encode(img_bytes).decode("utf-8"))
        doc.close()
        return images
    except Exception:
        return []


def is_probably_image(filename: str, mime: str = "") -> bool:
    m = (mime or "").lower()
    if m.startswith("image/"):
        return True
    return Path(filename).suffix.lower() in _IMAGE_EXT


def extract_text_from_bytes(filename: str, data: bytes, mime: str = "") -> str:
    """Return best-effort UTF-8 text; empty string if binary image (use vision path instead)."""
    if not data:
        return ""
    if is_probably_image(filename, mime):
        return ""

    ext = Path(filename).suffix.lower()
    name = Path(filename).name

    try:
        if ext == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            parts: list[str] = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)
            return "\n".join(parts).strip()

        if ext == ".docx":
            from docx import Document

            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text).strip()

        if ext == ".pptx":
            from pptx import Presentation

            prs = Presentation(io.BytesIO(data))
            texts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n".join(texts).strip()

        if ext in {".xlsx", ".xlsm"}:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            lines: list[str] = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        lines.append("\t".join(cells))
            wb.close()
            return "\n".join(lines).strip()

        if ext == ".csv":
            decoded = data.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(decoded))
            return "\n".join(",".join(row) for row in reader).strip()

        if ext in {
            ".txt",
            ".md",
            ".json",
            ".xml",
            ".html",
            ".htm",
            ".css",
            ".js",
            ".ts",
            ".tsx",
            ".jsx",
            ".py",
            ".java",
            ".cpp",
            ".c",
            ".h",
            ".cs",
            ".php",
            ".go",
            ".rs",
            ".rb",
            ".swift",
            ".kt",
            ".sql",
            ".yaml",
            ".yml",
            ".sh",
            ".r",
            ".m",
            ".scala",
            ".vue",
            ".svelte",
        }:
            return data.decode("utf-8", errors="replace").strip()

    except Exception:
        return f"[Could not extract text from {name}]"

    return data.decode("utf-8", errors="replace").strip()
